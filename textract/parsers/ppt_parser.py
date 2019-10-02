import pptx
import shutil
from tempfile import mkdtemp

from .utils import ShellParser


class Parser(ShellParser):
    """Extract text from ppt file using libreoffice and python-pptx
    """

    def extract(self, filename, **kwargs):
        # convert ppt to pptx with libre office
        temp_dir = mkdtemp()
        _ = self.convert_ppt(filename, temp_dir)

        # parse resulting pptx
        text = self.extract_pptx(temp_dir + "/" + filename + "x")

        shutil.rmtree(temp_dir)

        return text


    def convert_ppt(self, filename, directory):

        stdout, _ = self.run(["soffice", "--headless", "--convert-to", "pptx", filename, "--outdir", directory])

        return stdout


    def extract_pptx(self, filename):
        presentation = pptx.Presentation(filename)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return '\n\n'.join(text_runs)
